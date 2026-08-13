"""Generate and validate the editable 15-slide VideoMind Agent competition deck."""

from __future__ import annotations

import json
import subprocess
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR, MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt
from pptx.oxml.xmlchemy import OxmlElement


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "docs" / "competition"
PPTX = OUT / "VideoMind-Agent-参赛产品说明.pptx"
PDF = OUT / "VideoMind-Agent-参赛产品说明.pdf"
REPORT = OUT / "validation-report.json"
PDF_SCRIPT = Path(__file__).with_name("generate_competition_ppt.ps1")

SW, SH = 13.333333, 7.5
FONT = "Microsoft YaHei"
BG, DARK, GREEN = "F3F4EC", "244F3B", "3D6B53"
INK, MUTED, WHITE = "20251F", "6C746D", "FFFFFF"
PALE, LINE, ACCENT = "E4EBE2", "CAD5CB", "7B9B68"


def rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value)


def text(slide, value, x, y, w, h, size=18, color=INK, bold=False,
         align=PP_ALIGN.LEFT, name=None, valign=MSO_ANCHOR.TOP):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    if name:
        shape.name = name
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = 0
    frame.margin_top = frame.margin_bottom = 0
    frame.vertical_anchor = valign
    p = frame.paragraphs[0]
    p.text = value
    p.alignment = align
    p.font.name = FONT
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = rgb(color)
    for run in p.runs:
        run.font.name = FONT
    return shape


def box(slide, x, y, w, h, fill=WHITE, line=LINE, radius=True, name=None):
    kind = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    if name:
        shape.name = name
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(line)
    shape.line.width = Pt(1)
    return shape


def line(slide, x1, y1, x2, y2, color=ACCENT, arrow=True):
    shape = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    shape.line.color.rgb = rgb(color)
    shape.line.width = Pt(2)
    if arrow:
        # python-pptx exposes no public arrowhead API; set native DrawingML.
        ln = shape.line._get_or_add_ln()
        tail = OxmlElement("a:tailEnd")
        tail.set("type", "none")
        head = OxmlElement("a:headEnd")
        head.set("type", "triangle")
        head.set("w", "sm")
        head.set("len", "sm")
        ln.append(tail)
        ln.append(head)
    return shape


def title(slide, value, number, eyebrow="VIDEOMIND AGENT"):
    text(slide, eyebrow, .75, .38, 5.8, .25, 10, GREEN, True)
    text(slide, value, .75, .72, 11.8, .62, 29, DARK, True, name="SLIDE_TITLE")
    shape = box(slide, .75, 1.50, .68, .055, ACCENT, ACCENT, False)
    shape.line.fill.background()
    text(slide, f"{number:02d} / 15", 11.72, 6.94, .9, .22, 9, MUTED, False, PP_ALIGN.RIGHT)


def pill(slide, value, x, y, w, fill=PALE, color=DARK):
    box(slide, x, y, w, .39, fill, fill)
    text(slide, value, x, y + .08, w, .22, 11, color, True, PP_ALIGN.CENTER)


def card(slide, head, body, x, y, w, h, tag=None, fill=WHITE):
    box(slide, x, y, w, h, fill, LINE)
    if tag:
        text(slide, tag, x + .25, y + .22, w - .5, .2, 9, GREEN, True)
    hy = y + (.58 if tag else .3)
    text(slide, head, x + .25, hy, w - .5, .43, 18, DARK, True)
    text(slide, body, x + .25, hy + .62, w - .5, h - 1.1, 13, MUTED)


def placeholder(slide, label, x, y, w, h, page):
    shape = box(slide, x, y, w, h, PALE, ACCENT, False, f"SCREENSHOT_PLACEHOLDER_P{page}")
    shape.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    shape.line.width = Pt(1.5)
    text(slide, label, x + .25, y + h / 2 - .55, w - .5, .7, 16, GREEN, True,
         PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    text(slide, "在 PowerPoint / WPS 中删除此框并插入实际截图", x + .25, y + h / 2 + .35,
         w - .5, .28, 10, MUTED, False, PP_ALIGN.CENTER)


def step_flow(slide, items, x, y, total_w, h=.95, font=12):
    gap = .19
    bw = (total_w - gap * (len(items) - 1)) / len(items)
    # Connectors first, so they sit behind nodes.
    for i in range(len(items) - 1):
        x1 = x + i * (bw + gap) + bw
        line(slide, x1, y + h / 2, x1 + gap, y + h / 2)
    for i, item in enumerate(items):
        bx = x + i * (bw + gap)
        box(slide, bx, y, bw, h, WHITE if i % 2 == 0 else PALE, LINE)
        text(slide, item, bx + .08, y + .2, bw - .16, h - .35, font, DARK, True,
             PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)


def new_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = rgb(BG)
    return slide


def build() -> Presentation:
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(SW), Inches(SH)

    # 1 — cover
    s = new_slide(prs)
    pill(s, "迅雷校园 AI 产品创造营 · 创新 AI 工具", .75, .58, 3.72)
    text(s, "VideoMind Agent", .75, 1.45, 9.8, .82, 43, DARK, True, name="SLIDE_TITLE")
    text(s, "AI 视频理解与双语字幕智能工作台", .78, 2.48, 9.6, .45, 23, GREEN, True)
    text(s, "让长视频从“只能播放”\n变成“可识别、可校对、可翻译、可检索、可总结、可问答”的结构化知识载体。",
         .78, 3.42, 9.7, 1.05, 19, INK)
    box(s, 11.18, .95, 1.12, 4.55, DARK, DARK, False)
    text(s, "VIDEO\n↓\nKNOWLEDGE", 11.36, 1.85, .76, 2.5, 13, WHITE, True,
         PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    text(s, "参赛成员：________    学校：________    专业：________\nDemo：________________________    GitHub：________________________",
         .78, 6.02, 9.85, .68, 11, MUTED)
    text(s, "01 / 15", 11.72, 6.94, .9, .22, 9, MUTED, False, PP_ALIGN.RIGHT)

    # 2 — context
    s = new_slide(prs); title(s, "视频越来越多，信息却仍被困在线性时间轴里", 2)
    card(s, "定位效率低", "要找到一个观点，往往只能拖动进度条、反复试听。", .75, 2.08, 3.62, 2.9, "PAIN 01")
    card(s, "自动字幕不够准", "专业术语、口音与噪声会放大 ASR 错误，直接影响后续理解。", 4.86, 2.08, 3.62, 2.9, "PAIN 02")
    card(s, "工具链割裂", "转写、校对、翻译、摘要与问答分散在不同工具，上下文难复用。", 8.97, 2.08, 3.62, 2.9, "PAIN 03")
    text(s, "用户真正需要的不是另一块播放器，而是结构化、可搜索、可复用的信息。",
         2.05, 5.68, 9.2, .55, 22, DARK, True, PP_ALIGN.CENTER)

    # 3 — users
    s = new_slide(prs); title(s, "四类用户，共享同一个需求：更快获得可信信息", 3)
    users = [
        ("学生 / 学习者", "课程、讲座、公开课", "定位知识点慢", "可检索字幕与摘要"),
        ("内容创作者", "采访、素材、后期", "校对与翻译耗时", "可控字幕与 SRT"),
        ("研究 / 知识工作者", "会议、访谈、资料", "长视频难复用", "结构化理解与问答"),
        ("跨语言视频用户", "海外课程与内容", "语言成为门槛", "先校正、再翻译"),
    ]
    for i, (who, scene, pain, need) in enumerate(users):
        x = .75 + i * 3.0
        box(s, x, 2.05, 2.75, 4.0, WHITE, LINE)
        text(s, who, x + .22, 2.35, 2.3, .48, 17, DARK, True)
        pill(s, scene, x + .22, 3.0, 2.3)
        text(s, "痛点", x + .22, 3.78, .7, .22, 10, GREEN, True)
        text(s, pain, x + .22, 4.1, 2.3, .42, 13)
        text(s, "需求", x + .22, 4.88, .7, .22, 10, GREEN, True)
        text(s, need, x + .22, 5.2, 2.3, .46, 13)

    # 4 — positioning
    s = new_slide(prs); title(s, "时间戳字幕，是连接视频与 AI 理解的核心数据层", 4)
    line(s, 4.15, 2.85, 5.15, 2.85); line(s, 8.2, 2.85, 9.2, 2.85)
    box(s, 1.25, 2.3, 2.9, 1.1); text(s, "视频", 1.25, 2.62, 2.9, .35, 22, DARK, True, PP_ALIGN.CENTER)
    box(s, 5.15, 2.02, 3.05, 1.65, DARK, DARK)
    text(s, "时间戳字幕", 5.15, 2.42, 3.05, .4, 25, WHITE, True, PP_ALIGN.CENTER)
    text(s, "统一、可编辑、可追溯", 5.15, 3.02, 3.05, .25, 11, WHITE, False, PP_ALIGN.CENTER)
    box(s, 9.2, 2.3, 2.65, 1.1); text(s, "AI 理解", 9.2, 2.62, 2.65, .35, 22, DARK, True, PP_ALIGN.CENTER)
    card(s, "更快", "搜索、跳转、摘要和问答，减少重复观看。", 1.25, 4.65, 3.0, 1.5)
    card(s, "更准", "AI 校正 + 人工确认，降低错误传播。", 5.15, 4.65, 3.0, 1.5)
    card(s, "更可控", "人工编辑、处理记录与 fallback 保留最终决定权。", 9.05, 4.65, 3.0, 1.5)

    # 5 — workflow
    s = new_slide(prs); title(s, "一条连续工作流，把原始视频变成可复用知识", 5)
    step_flow(s, ["视频输入", "Whisper\n语音识别", "Transcript\nCorrection", "人工字幕\n编辑", "双语字幕", "AI 视频\n摘要", "视频问答"], .58, 2.8, 12.18, 1.15, 12)
    text(s, "识别", 2.2, 4.42, 1, .25, 10, GREEN, True, PP_ALIGN.CENTER)
    text(s, "校正与控制", 5.15, 4.42, 2, .25, 10, GREEN, True, PP_ALIGN.CENTER)
    text(s, "理解与复用", 9.45, 4.42, 2, .25, 10, GREEN, True, PP_ALIGN.CENTER)
    text(s, "同一份结构化字幕贯穿全流程，避免在多个工具间重复导入与丢失上下文。",
         1.3, 5.45, 10.7, .5, 18, DARK, True, PP_ALIGN.CENTER)

    # 6 — layers
    s = new_slide(prs); title(s, "三层架构，让交互、字幕质量与 AI 理解各司其职", 6)
    layers = [
        ("交互层", "视频播放  ·  倍速控制  ·  时间戳跳转  ·  播放同步", WHITE),
        ("字幕层", "Whisper 转写  ·  AI 字幕校正  ·  人工编辑  ·  双语字幕  ·  SRT 导出", PALE),
        ("AI 理解层", "视频摘要  ·  视频问答  ·  内容结构化", WHITE),
    ]
    for i, (name, body, fill) in enumerate(layers):
        y = 2.0 + i * 1.5
        if i < 2: line(s, 6.66, y + 1.12, 6.66, y + 1.5)
        box(s, 1.5, y, 10.33, 1.12, fill, LINE)
        text(s, name, 1.85, y + .34, 2.05, .38, 20, DARK, True)
        text(s, body, 4.12, y + .36, 7.2, .35, 14, INK)

    # 7 — correction stability
    s = new_slide(prs); title(s, "严格校验与可恢复机制，提高 TranscriptCorrection 稳定性", 7)
    step_flow(s, ["Whisper\n原始字幕", "字幕\n批处理", "LLM\nCorrection", "JSON\n校验", "自动\nRetry", "Fallback", "最终字幕"], .58, 2.15, 12.18, 1.05, 12)
    features = ["JSON 容错解析", "invalid output 自动重试", "失败批次 fallback", "retry / fallback metadata", "Whisper 时间戳保护"]
    coords = [(1, 4.08), (4.72, 4.08), (8.44, 4.08), (2.86, 5.03), (6.58, 5.03)]
    for i, value in enumerate(features):
        x, y = coords[i]; pill(s, value, x, y, 3.15, DARK if i == 4 else PALE, WHITE if i == 4 else DARK)
    text(s, "安全约束：不新增 cue、不丢失 cue、不猜测 ID；start / end 始终沿用 Whisper 原始时间轴。",
         1.05, 6.12, 11.2, .4, 15, DARK, True, PP_ALIGN.CENTER)

    # 8 — human in loop
    s = new_slide(prs); title(s, "Human-in-the-loop：AI 提效，人对最终质量负责", 8)
    placeholder(s, "【替换为 SubtitleEditor 实际截图】", 5.42, 1.95, 7.2, 4.55, 8)
    text(s, "编辑闭环", .75, 2.05, 3.3, .38, 20, DARK, True)
    for i, value in enumerate(["时间戳定位", "AI 校正结果", "人工修改", "单条 / 批量保存", "处理记录", "SRT 导出"]):
        text(s, f"✓  {value}", .8, 2.75 + i * .58, 3.75, .35, 15, INK)
    text(s, "AI baseline、已保存编辑、未保存草稿三层状态清晰分离。", .8, 6.3, 3.9, .5, 12, MUTED)

    # 9 — realtime sync
    s = new_slide(prs); title(s, "同一个媒体时钟，驱动播放器字幕与编辑器高亮", 9)
    placeholder(s, "【替换为播放器 + SubtitleEditor 同步截图】", 8.2, 1.95, 4.42, 4.55, 9)
    box(s, .75, 2.0, 6.7, 3.72, WHITE, LINE)
    text(s, "requestAnimationFrame 播放时钟", 1.05, 2.33, 6.1, .35, 18, DARK, True, PP_ALIGN.CENTER)
    box(s, 1.45, 3.15, 2.15, .65, PALE); text(s, "video.currentTime", 1.45, 3.35, 2.15, .25, 13, DARK, True, PP_ALIGN.CENTER)
    line(s, 3.6, 3.48, 4.5, 3.48)
    box(s, 4.5, 3.15, 2.15, .65, DARK, DARK); text(s, "findActiveCueId", 4.5, 3.35, 2.15, .25, 13, WHITE, True, PP_ALIGN.CENTER)
    line(s, 5.58, 3.8, 4.32, 4.55); line(s, 5.58, 3.8, 6.4, 4.55)
    pill(s, "SubtitleTrack", 3.12, 4.62, 1.9); pill(s, "SubtitleEditor", 5.25, 4.62, 1.9)
    text(s, "seek 即时同步  ·  active cue 实时更新  ·  自动滚动与高亮解耦", .82, 6.15, 6.55, .3, 13, GREEN, True, PP_ALIGN.CENTER)

    # 10 — follow UX
    s = new_slide(prs); title(s, "高亮实时变化，滚动始终尊重人工操作", 10)
    box(s, .75, 2.05, 5.95, 4.25, WHITE, LINE)
    text(s, "自动跟随", 1.15, 2.42, 2.2, .38, 20, DARK, True)
    text(s, "当前 cue 定位到编辑区域约 30% 高度：\n上方保留已播放内容，下方展示更多即将播放字幕。", 1.15, 3.05, 4.85, .85, 15)
    box(s, 1.38, 4.4, 4.5, 1.2, PALE, LINE)
    text(s, "上一条\n▶ 当前播放字幕\n下一条 · 下一条 · 下一条", 1.65, 4.55, 3.95, .85, 13, DARK, True, PP_ALIGN.CENTER)
    text(s, "以下操作优先于自动跟随", 7.55, 2.2, 4.4, .38, 20, DARK, True)
    for i, value in enumerate(["滚动字幕列表", "拖动页面滚动条", "搜索字幕", "编辑 textarea"]):
        pill(s, value, 7.55, 3.05 + i * .72, 3.9)
    text(s, "停止人工操作约 2 秒后恢复跟随；\n暂停滚动不影响 active cue 实时高亮。", 7.55, 6.12, 4.05, .52, 13, MUTED)

    # 11 — bilingual
    s = new_slide(prs); title(s, "先校正，再翻译，避免错误从原文传播到译文", 11)
    step_flow(s, ["原始字幕", "AI 校正", "AI 翻译", "复用时间戳", "双语字幕"], .75, 2.12, 7.1, .92, 12)
    text(s, "时间轴只维护一份，原文与译文共享 cue 边界。\n人工编辑后的 effective text 继续用于导出、摘要与问答。", .8, 3.85, 6.55, .82, 15)
    pill(s, "SOURCE", .8, 5.35, 1.5, DARK, WHITE); pill(s, "TRANSLATION", 2.52, 5.35, 2.0)
    placeholder(s, "【替换为双语字幕播放器截图】", 8.35, 1.95, 4.27, 4.55, 11)

    # 12 — summary and Q&A
    s = new_slide(prs); title(s, "从“按时间找信息”，升级为“按问题找信息”", 12)
    box(s, .75, 1.95, 5.7, 4.55, WHITE, LINE)
    text(s, "AI 视频摘要", 1.1, 2.32, 4.8, .45, 22, DARK, True)
    text(s, "提炼主题、关键结论与内容章节，\n帮助用户先理解全貌，再决定观看位置。", 1.1, 3.0, 4.8, .72, 15)
    placeholder(s, "【替换为 Summary 实际截图】", 1.1, 4.1, 5.0, 1.72, 12)
    box(s, 6.88, 1.95, 5.7, 4.55, WHITE, LINE)
    text(s, "AI 视频问答", 7.23, 2.32, 4.8, .45, 22, DARK, True)
    text(s, "直接围绕视频内容提出问题，\n新问题使用最新人工字幕作为上下文。", 7.23, 3.0, 4.8, .72, 15)
    placeholder(s, "【替换为 Q&A 实际截图】", 7.23, 4.1, 5.0, 1.72, 12)

    # 13 — native AI architecture
    s = new_slide(prs); title(s, "AI 原生能力建立在结构化字幕与严格工程约束之上", 13)
    step_flow(s, ["Video", "Whisper", "Structured\nSubtitle", "LLM", "Correction / Translation / Summary / QA"], .75, 2.2, 11.82, 1.08, 12)
    tech = ["Whisper", "LLM", "结构化字幕", "Human-in-the-loop", "时间轴同步", "JSON 容错"]
    for i, value in enumerate(tech):
        x = 1.05 + (i % 3) * 3.9; y = 4.32 + (i // 3) * .9
        pill(s, value, x, y, 3.3, DARK if i == 3 else PALE, WHITE if i == 3 else DARK)
    text(s, "当前架构围绕本地视频、Whisper、结构化字幕文件与 LLM API 展开。",
         1.15, 6.28, 11.0, .35, 13, MUTED, False, PP_ALIGN.CENTER)

    # 14 — feasibility
    s = new_slide(prs); title(s, "核心链路已可运行，商业化仍属于下一阶段规划", 14)
    box(s, .75, 1.95, 5.95, 4.6, WHITE, LINE)
    text(s, "当前完成情况", 1.12, 2.3, 5.1, .38, 20, DARK, True)
    text(s, "✓ 视频播放     ✓ 字幕生成\n✓ AI 校正      ✓ 字幕编辑\n✓ 双语字幕     ✓ AI 摘要\n✓ 视频问答     ✓ SRT 导出", 1.15, 3.02, 4.75, 1.38, 14)
    text(s, "Backend", 1.15, 4.83, 1.0, .22, 10, GREEN, True); text(s, "89 passed · 27 subtests passed", 2.25, 4.8, 3.7, .28, 13, DARK, True)
    text(s, "Frontend", 1.15, 5.3, 1.0, .22, 10, GREEN, True); text(s, "41 passed · 0 failed", 2.25, 5.27, 3.7, .28, 13, DARK, True)
    text(s, "Build", 1.15, 5.77, 1.0, .22, 10, GREEN, True); text(s, "Vite build passed", 2.25, 5.74, 3.7, .28, 13, DARK, True)
    box(s, 7.15, 1.95, 5.43, 4.6, PALE, LINE)
    text(s, "商业化方向", 7.52, 2.3, 4.65, .38, 20, DARK, True)
    card(s, "个人版", "学习、创作与个人知识整理", 7.52, 3.05, 1.4, 2.0)
    card(s, "专业版", "高频字幕与跨语言工作流", 9.15, 3.05, 1.4, 2.0)
    card(s, "团队 / API", "协作、批处理与能力集成", 10.78, 3.05, 1.4, 2.0)
    text(s, "以上为未来规划，并非当前已上线收费能力。", 7.52, 5.8, 4.65, .3, 11, MUTED, False, PP_ALIGN.CENTER)

    # 15 — roadmap and ending
    s = new_slide(prs); title(s, "从稳定可用，走向视频知识化", 15)
    roads = [
        ("阶段 1", "稳定可用", "字幕、编辑、同步、导出"),
        ("阶段 2", "理解增强", "摘要定位、问答证据、术语词表"),
        ("阶段 3", "视频知识化", "多语言、多视频检索、知识库、协作与 API"),
    ]
    for i in range(2): line(s, 4.45 + i * 4.08, 2.93, 4.78 + i * 4.08, 2.93)
    for i, (stage, head, body) in enumerate(roads):
        x = .75 + i * 4.08; fill = DARK if i == 0 else WHITE
        box(s, x, 2.0, 3.75, 1.95, fill, DARK if i == 0 else LINE)
        text(s, stage, x + .28, 2.28, 3.15, .2, 10, WHITE if i == 0 else GREEN, True)
        text(s, head, x + .28, 2.7, 3.15, .38, 20, WHITE if i == 0 else DARK, True)
        text(s, body, x + .28, 3.25, 3.15, .42, 12, WHITE if i == 0 else MUTED)
    text(s, "VideoMind Agent", .75, 4.65, 6.8, .55, 30, DARK, True)
    text(s, "让 AI 不只是“看见视频”，\n而是真正理解、整理和复用视频中的知识。", .75, 5.35, 7.6, .78, 19)
    text(s, "Demo：________________    GitHub：________________\n团队：________________    联系方式：________________",
         8.2, 5.1, 4.4, .9, 11, MUTED)

    return prs


def validate(path: Path) -> dict:
    prs = Presentation(path)
    errors, warnings = [], []
    if len(prs.slides) != 15:
        errors.append(f"slide_count={len(prs.slides)}, expected=15")
    ratio = prs.slide_width / prs.slide_height
    if abs(ratio - 16 / 9) > 0.002:
        errors.append(f"aspect_ratio={ratio:.6f}, expected={16/9:.6f}")
    editable_text_shapes = 0
    picture_shapes = 0
    placeholders = set()
    slide_checks = []
    for number, slide in enumerate(prs.slides, 1):
        title_shapes = [sh for sh in slide.shapes if sh.name == "SLIDE_TITLE" and sh.has_text_frame and sh.text.strip()]
        if not title_shapes:
            errors.append(f"slide {number}: missing title")
        out_of_bounds = []
        for sh in slide.shapes:
            if sh.left < 0 or sh.top < 0 or sh.left + sh.width > prs.slide_width or sh.top + sh.height > prs.slide_height:
                out_of_bounds.append(sh.name)
            if sh.has_text_frame and sh.text.strip():
                editable_text_shapes += 1
            if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
                picture_shapes += 1
            if sh.name.startswith("SCREENSHOT_PLACEHOLDER_"):
                placeholders.add(number)
        if out_of_bounds:
            errors.append(f"slide {number}: objects outside boundary: {out_of_bounds}")
        slide_checks.append({"slide": number, "title": title_shapes[0].text if title_shapes else None,
                             "shape_count": len(slide.shapes), "out_of_bounds": out_of_bounds})
    required_placeholders = {8, 9, 11, 12}
    if not required_placeholders.issubset(placeholders):
        errors.append(f"missing screenshot placeholders: {sorted(required_placeholders - placeholders)}")
    if picture_shapes:
        warnings.append(f"picture_shapes={picture_shapes}; deck must not contain full-slide raster pages")
    report = {
        "valid": not errors,
        "slide_count": len(prs.slides),
        "size_inches": [round(prs.slide_width / 914400, 3), round(prs.slide_height / 914400, 3)],
        "aspect_ratio": round(ratio, 6),
        "editable_text_shapes": editable_text_shapes,
        "picture_shapes": picture_shapes,
        "screenshot_placeholder_pages": sorted(placeholders),
        "errors": errors,
        "warnings": warnings,
        "slides": slide_checks,
    }
    REPORT.write_text(json.dumps(report, ensure_ascii=False, indent=2), encoding="utf-8")
    if errors:
        raise RuntimeError("PPTX validation failed: " + "; ".join(errors))
    return report


def export_pdf() -> bool:
    try:
        subprocess.run([
            "powershell.exe", "-NoProfile", "-ExecutionPolicy", "Bypass", "-File", str(PDF_SCRIPT),
            "-PptxPath", str(PPTX), "-PdfPath", str(PDF)
        ], check=True, timeout=180)
        return PDF.exists() and PDF.stat().st_size > 0
    except Exception as exc:
        print(f"PDF export failed (PPTX remains valid): {exc}")
        return False


def main() -> None:
    OUT.mkdir(parents=True, exist_ok=True)
    prs = build()
    prs.save(PPTX)
    report = validate(PPTX)
    print(f"PPTX generated and validated: {PPTX}")
    print(json.dumps({k: report[k] for k in ("slide_count", "size_inches", "aspect_ratio", "editable_text_shapes", "picture_shapes", "screenshot_placeholder_pages")}, ensure_ascii=False))
    print(f"PDF exported: {export_pdf()}")


if __name__ == "__main__":
    main()
